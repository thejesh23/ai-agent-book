#!/usr/bin/env python3
"""
Experiment 2-6: Generating Presentations from Papers Using Agent Skills (Self-built Isomorphic Skills Mechanism)

This demo reproduces the idea from Chapter 2 "Agent Skills / Progressive Disclosure" of "Deep Understanding of AI Agents".
Since the Anthropic key is invalid, we use OpenAI (gpt-5.6-luna) + a self-built mechanism isomorphic to Anthropic
Skills to demonstrate, with the core concept being "Progressive Disclosure":

  Layer 1 (Metadata): The system prompt at Agent startup only contains each Skill's name +
                    description (thin catalog, hundreds of tokens), without specific procedures.
  Layer 2 (Core Process): When needed by the task, the Agent actively loads the complete SKILL.md using the read_skill tool.
  Layer 3 (Details): The Agent can further use read_skill_file to read reference.md / script source code.

Then the Agent uses the bundled script scripts/generate_pptx.py (via the run_skill_script tool) with
python-pptx to generate a real .pptx, and reads back to verify the page count and each page title.

Run:
    export OPENAI_API_KEY=sk-...
    python demo.py
"""

import argparse
import json
import os
import sys
from pathlib import Path

from openai import OpenAI
from pptx import Presentation

# Read OPENAI_API_KEY from .env in the same directory (if python-dotenv is installed)
try:
    from dotenv import load_dotenv

    load_dotenv(Path(__file__).resolve().parent / ".env")
except ImportError:
    pass

# ---------------------------------------------------------------------------
# Paths and Configuration
# ---------------------------------------------------------------------------
ROOT = Path(__file__).resolve().parent
SKILLS_DIR = ROOT / "skills"
PAPER_PATH = ROOT / "papers" / "sample_paper.md"
OUTPUT_DIR = ROOT / "output"
MODEL = os.environ.get("OPENAI_MODEL", "gpt-5.6-luna")


def log(msg: str) -> None:
    print(msg, flush=True)


# ---------------------------------------------------------------------------
# Layer 1: Scan the skills/ directory at startup, only read the frontmatter of each SKILL.md
# (name + description), compose a thin catalog and inject it into the system prompt. This step deliberately "only looks at the catalog".
# ---------------------------------------------------------------------------
def parse_frontmatter(skill_md: str) -> dict:
    """ Parse name / description from the --- YAML frontmatter --- at the top of SKILL.md."""
    meta = {}
    if not skill_md.startswith("---"):
        return meta
    end = skill_md.find("---", 3)
    if end == -1:
        return meta
    for line in skill_md[3:end].splitlines():
        if ":" in line:
            k, v = line.split(":", 1)
            meta[k.strip()] = v.strip()
    return meta


def scan_skill_catalog() -> dict:
    """ Return {skill_name: {"description":..., "dir": Path}}, containing only metadata."""
    catalog = {}
    for skill_md in sorted(SKILLS_DIR.glob("*/SKILL.md")):
        meta = parse_frontmatter(skill_md.read_text(encoding="utf-8"))
        name = meta.get("name") or skill_md.parent.name
        catalog[name] = {
            "description": meta.get("description", ""),
            "dir": skill_md.parent,
        }
    return catalog


def build_system_prompt(catalog: dict) -> str:
    lines = [
        "You are an assistant capable of using Agent Skills. You do not know the detailed process of each Skill in advance;",
        "you only see a \"thin catalog\" below—the name and description (routing conditions) of each Skill.",
        "",
        "When a task requires a certain Skill, you must:",
        "  1) First use read_skill(name) to load its complete SKILL.md (Layer 2: core process);",
        "  2) If implementation/style details are needed, use read_skill_file(name, path) to read subdocuments or scripts (Layer 3);",
        "  3) Follow the conventions in SKILL.md to call the bundled script via run_skill_script to complete the task.",
        "Do not guess the invocation method of a Skill without using read_skill.",
        "",
        "== Installed Skills (Thin Catalog, Metadata Only) ==",
    ]
    for name, info in catalog.items():
        lines.append(f"- {name}: {info['description']}")
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# Tool implementations: read_skill / read_skill_file / run_skill_script
# These are the channels for "Progressive Disclosure"—Layer 2 and Layer 3 content only enters the context when invoked.
# ---------------------------------------------------------------------------
def tool_read_skill(catalog: dict, name: str) -> str:
    info = catalog.get(name)
    if not info:
        return f"[error] Skill not found: {name}"
    content = (info["dir"] / "SKILL.md").read_text(encoding="utf-8")
    log(f"\n  >>> [Progressive Disclosure·Layer 2] Agent calls read_skill('{name}')，"
        f"Loaded complete SKILL.md ({len(content)} characters)")
    return content


def tool_read_skill_file(catalog: dict, name: str, rel_path: str) -> str:
    info = catalog.get(name)
    if not info:
        return f"[error] Skill not found: {name}"
    target = (info["dir"] / rel_path).resolve()
    # Anti-directory-traversal: must be within the skill directory
    if not str(target).startswith(str(info["dir"].resolve())):
        return f"[error] Illegal path: {rel_path}"
    if not target.exists():
        return f"[error] File does not exist: {rel_path}"
    content = target.read_text(encoding="utf-8")
    log(f"  >>> [Progressive Disclosure·Layer 3] Agent calls read_skill_file('{name}', '{rel_path}')，"
        f"Loaded subdocument ({len(content)} characters)")
    return content


def tool_run_skill_script(catalog: dict, name: str, script: str, payload: str,
                          out_path: Path) -> str:
    info = catalog.get(name)
    if not info:
        return f"[error] Skill not found: {name}"
    script_path = (info["dir"] / "scripts" / script).resolve()
    if not script_path.exists():
        return f"[error] Script does not exist: {script}"

    # Dynamically load bundled script (it is part of the Skill)
    import importlib.util
    spec = importlib.util.spec_from_file_location("bundled_generator", script_path)
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)

    try:
        data = json.loads(payload) if isinstance(payload, str) else payload
    except json.JSONDecodeError as e:
        return f"[error] payload is not valid JSON: {e}"

    log(f"  >>> [Executing bundled script] run_skill_script('{name}', '{script}') "
        f"Generate {out_path.name} ...")
    result = module.build_presentation(data, str(out_path))
    return json.dumps(result, ensure_ascii=False)


TOOLS = [
    {
        "type": "function",
        "function": {
            "name": "read_skill",
            "description": "Load the complete SKILL.md of the specified Skill (core process, progressive disclosure of the second layer).",
            "parameters": {
                "type": "object",
                "properties": {"name": {"type": "string", "description": "Skill name"}},
                "required": ["name"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "read_skill_file",
            "description": "Read the sub-documents or script source code within a certain Skill directory (details, progressive disclosure layer 3).",
            "parameters": {
                "type": "object",
                "properties": {
                    "name": {"type": "string"},
                    "path": {"type": "string", "description": "Path relative to the skill directory, such as reference.md or scripts/generate_pptx.py"},
                },
                "required": ["name", "path"],
            },
        },
    },
    {
        "type": "function",
        "function": {
            "name": "run_skill_script",
            "description": "Execute the script bound to a Skill to produce actual output (e.g., generate a pptx).",
            "parameters": {
                "type": "object",
                "properties": {
                    "name": {"type": "string"},
                    "script": {"type": "string", "description": "Script file name, e.g., generate_pptx.py"},
                    "payload": {"type": "string", "description": "JSON string passed to the script (outline)"},
                },
                "required": ["name", "script", "payload"],
            },
        },
    },
]


def dispatch(catalog: dict, name: str, args: dict, out_path: Path) -> str:
    if name == "read_skill":
        return tool_read_skill(catalog, args["name"])
    if name == "read_skill_file":
        return tool_read_skill_file(catalog, args["name"], args["path"])
    if name == "run_skill_script":
        return tool_run_skill_script(catalog, args["name"], args["script"],
                                     args["payload"], out_path)
    return f"[error] Unknown tool: {name}"


# ---------------------------------------------------------------------------
# Main flow: agentic loop
# ---------------------------------------------------------------------------
def run_agent(paper_path: Path, model: str, out_path: Path,
              max_turns: int = 8) -> Path | None:
    # If OPENAI_API_KEY exists, connect directly to official; otherwise fallback to OPENROUTER_API_KEY
    # (gpt-* model names will be mapped to openai/...). If neither is present, give a clear error.
    from openrouter_fallback import resolve_llm

    if not os.environ.get("OPENAI_API_KEY") and not os.environ.get("OPENROUTER_API_KEY"):
        log("Error: OPENAI_API_KEY is not set, nor is OPENROUTER_API_KEY (general fallback).")
        log("Please export OPENAI_API_KEY=sk-... or export OPENROUTER_API_KEY=sk-or-...")
        log("(When there is no key, you can use --offline to follow the built-in outline, deterministically reproduce three-layer progressive disclosure, and generate pptx.)")
        sys.exit(1)

    api_key, base_url, model = resolve_llm(
        model=model,
        primary_keys=("OPENAI_API_KEY",),
        primary_base_url=os.getenv("OPENAI_BASE_URL") or None,
    )
    # timeout + automatic retry: a single network/SSL glitch should not crash the entire agentic loop
    client = OpenAI(api_key=api_key, base_url=base_url, timeout=60.0, max_retries=3)
    catalog = scan_skill_catalog()

    system_prompt = build_system_prompt(catalog)
    log("=" * 72)
    log("[Layer 1 · Metadata] When the Agent starts, it only sees this thin Skill directory (system prompt):")
    log("-" * 72)
    log(system_prompt)
    log("-" * 72)
    log(f"(thin directory approx {len(system_prompt)}  characters / hundreds of tokens; the detailed flow of each Skill is not in the context at the moment)")
    log("=" * 72)

    paper = paper_path.read_text(encoding="utf-8")
    user_task = (
        "Please turn the following paper into an 8-12 page presentation (including title page, table of contents, problem background,"
        "Method overview, key results, limitations, summary page), the total number of pages must be between 8 and 12."
        "First determine which Skill to use, then strictly follow the page order and constraints of its SKILL.md.\n\n"
        "=== Full Paper ===\n" + paper
    )

    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": user_task},
    ]

    log("\n[Task Assignment] The agent is required to generate a presentation from a paper. Observe how it progressively discloses information on demand: \n")

    final_result = None
    for turn in range(1, max_turns + 1):
        resp = client.chat.completions.create(
            model=model,
            messages=messages,
            tools=TOOLS,
            temperature=0.2,
        )
        msg = resp.choices[0].message
        messages.append(msg.model_dump(exclude_none=True))

        if not msg.tool_calls:
            log(f"\n【Agent No. {turn}  Round·Conclusion】\n{msg.content}")
            break

        for tc in msg.tool_calls:
            fn = tc.function.name
            try:
                args = json.loads(tc.function.arguments or "{}")
            except json.JSONDecodeError:
                args = {}
            log(f"\n[Agent Round {turn} ] call tool -> {fn}({', '.join(f'{k}={_short(v)}' for k, v in args.items())})")
            result = dispatch(catalog, fn, args, out_path)
            if fn == "run_skill_script" and not result.startswith("[error]"):
                final_result = json.loads(result)
                log(f"  >>> Generation result:{result}")
            messages.append({
                "role": "tool",
                "tool_call_id": tc.id,
                "content": result,
            })

    if final_result:
        return Path(final_result["path"])
    return None


# ---------------------------------------------------------------------------
# Offline reproduction: When there is no OpenAI key, use the built-in outline (papers/sample_outline.json) deterministically
# Walk through the exact same three-level progressive disclosure and tool channel as online (read_skill / read_skill_file /
# run_skill_script), so that pptx can be genuinely generated and verified without any API access.
# The only difference is that "which Skill to use and what to write in the outline" is given by a preset script, not decided by the model in real time.
# ---------------------------------------------------------------------------
OUTLINE_PATH = ROOT / "papers" / "sample_outline.json"


def run_offline(out_path: Path) -> Path | None:
    catalog = scan_skill_catalog()
    system_prompt = build_system_prompt(catalog)
    log("=" * 72)
    log(" [Offline Mode] Does not call OpenAI, uses built-in outline to deterministically reproduce the three-level progressive disclosure.")
    log(" [Layer 1 · Metadata] On startup, only this thin Skill directory (system prompt) is visible:")
    log("-" * 72)
    log(system_prompt)
    log("-" * 72)
    log(f"(thin directory approx {len(system_prompt)} characters; the detailed flow of each Skill is not in the context at this moment)")
    log("=" * 72)

    if not OUTLINE_PATH.exists():
        log(f" Error: Built-in outline does not exist:{OUTLINE_PATH}")
        return None

    # Same tool channel as the online agentic loop, except the call sequence is given by the script
    log("\n[Offline Replay] According to SKILL.md convention, load and call bundled scripts layer by layer:")
    dispatch(catalog, "read_skill", {"name": "pptx"}, out_path)
    dispatch(catalog, "read_skill_file",
             {"name": "pptx", "path": "reference.md"}, out_path)

    payload = OUTLINE_PATH.read_text(encoding="utf-8")
    result = dispatch(catalog, "run_skill_script",
                      {"name": "pptx", "script": "generate_pptx.py", "payload": payload},
                      out_path)
    if result.startswith("[error]"):
        log(f"  >>> Generation failed:{result}")
        return None
    log(f"  >>> Generation result:{result}")
    return Path(json.loads(result)["path"])


def _short(v, n=48):
    s = str(v).replace("\n", " ")
    return s if len(s) <= n else s[:n] + "…"


# ---------------------------------------------------------------------------
# Verification: Use python-pptx to reopen the generated file, read back the page count and each page title, proving it is a valid pptx.
# ---------------------------------------------------------------------------
def verify_pptx(path: Path) -> None:
    log("\n" + "=" * 72)
    log(" [Verification] Use python-pptx to reopen the generated file, read back the page count and each page title:")
    log("-" * 72)
    prs = Presentation(str(path))
    slides = list(prs.slides)
    log(f" File: {path}")
    log(f" Total pages: {len(slides)}")
    for i, slide in enumerate(slides, 1):
        first_text = "(empty)"
        for shp in slide.shapes:
            if shp.has_text_frame and shp.text_frame.text.strip():
                first_text = shp.text_frame.text.strip().splitlines()[0]
                break
        log(f"  Page {i:>2} title: {first_text}")
    log("-" * 72)
    log(f" Verification passed: This is a valid .pptx that can be opened by python-pptx / PowerPoint ({len(slides)} pages).")
    log("=" * 72)


def parse_args():
    p = argparse.ArgumentParser(
        description=" Experiment 2-6: Use Agent Skills' \"progressive disclosure\" to generate a presentation from a paper."
                    " The Agent only sees a thin Skill directory on startup, loads the pptx Skill's flow and scripts layer by layer as needed,"
                    " then uses python-pptx to generate and verify output/presentation.pptx.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    p.add_argument("--paper", default=str(PAPER_PATH),
                   help=" Input paper/outline (markdown) path, default papers/sample_paper.md.")
    p.add_argument("--output", "-o", default=str(OUTPUT_DIR / "presentation.pptx"),
                   help=" Output .pptx path, default output/presentation.pptx.")
    p.add_argument("--model", default=MODEL,
                   help=" OpenAI model name, defaults to environment variable OPENAI_MODEL, otherwise gpt-5.6-luna.")
    p.add_argument("--max-turns", type=int, default=8,
                   help=" Maximum number of rounds for the agentic loop, default 8.")
    p.add_argument("--offline", action="store_true",
                   help=" Offline demo: Does not call OpenAI, uses built-in outline (papers/sample_outline.json)"
                        " to deterministically walk through the three-level progressive disclosure and generate pptx (no API key needed, reproducible).")
    return p.parse_args()


def main():
    args = parse_args()
    paper_path = Path(args.paper)
    out_path = Path(args.output)
    out_path.parent.mkdir(parents=True, exist_ok=True)

    if args.offline:
        pptx_path = run_offline(out_path)
    else:
        if not paper_path.exists():
            log(f" Error: Paper file does not exist:{paper_path}")
            sys.exit(1)
        pptx_path = run_agent(paper_path, args.model, out_path, args.max_turns)

    if pptx_path and pptx_path.exists():
        verify_pptx(pptx_path)
    else:
        log("\nNo pptx generated. Please check the logs above.")
        sys.exit(2)


if __name__ == "__main__":
    main()
