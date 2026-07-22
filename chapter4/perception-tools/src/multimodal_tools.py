"""
Multimodal understanding tools: web, documents, images, and videos.
"""
import json
import logging
import os
import traceback
from pathlib import Path
from typing import Optional, Union
import base64

import requests
from bs4 import BeautifulSoup
from dotenv import load_dotenv
from mcp.types import TextContent
from pydantic import Field

from base import ActionResponse, validate_file_path, download_file_from_url, is_url


load_dotenv()


async def read_webpage(
    url: str,
    extract_text: bool = True,
    extract_links: bool = False
) -> Union[str, TextContent]:
    """
    Read and extract content from a webpage.
    
    Args:
        url: URL of the webpage
        extract_text: Whether to extract main text content
        extract_links: Whether to extract all links
        
    Returns:
        TextContent with extracted webpage content
    """
    try:
        logging.info(f"📄 Reading webpage: {url}")
        
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"
        }
        
        response = requests.get(url, headers=headers, timeout=30)
        response.raise_for_status()
        
        soup = BeautifulSoup(response.content, 'html.parser')
        
        result = {
            "url": url,
            "title": soup.title.string if soup.title else "No title"
        }
        
        if extract_text:
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            text = soup.get_text()
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = ' '.join(chunk for chunk in chunks if chunk)
            
            result["text"] = text[:5000]  # Limit to first 5000 chars
            result["text_length"] = len(text)
        
        if extract_links:
            links = []
            for link in soup.find_all('a', href=True):
                links.append({
                    "text": link.get_text().strip(),
                    "href": link['href']
                })
            result["links"] = links[:50]  # Limit to first 50 links
        
        logging.info(f"✅ Successfully extracted webpage content")
        
        action_response = ActionResponse(
            success=True,
            message=result,
            metadata={"url": url}
        )
        
        return TextContent(
            type="text",
            text=json.dumps(action_response.model_dump())
        )
        
    except Exception as e:
        error_msg = f"Webpage reading failed: {str(e)}"
        logging.error(f"Webpage error: {traceback.format_exc()}")
        
        action_response = ActionResponse(
            success=False,
            message=error_msg,
            metadata={"error_type": "webpage_error", "url": url}
        )
        
        return TextContent(
            type="text",
            text=json.dumps(action_response.model_dump())
        )


def _sniff_document_type(path: Path) -> Optional[str]:
    """Detect .pdf/.docx/.pptx from magic bytes when the extension is unusable.

    A URL's path often carries no real extension (e.g.
    https://arxiv.org/pdf/2301.07041 -> ".07041"), so the downloaded temp
    file's suffix cannot be trusted to identify the format.
    """
    try:
        with open(path, 'rb') as f:
            header = f.read(4)
        if header.startswith(b'%PDF'):
            return '.pdf'
        if header.startswith(b'PK\x03\x04'):
            import zipfile
            with zipfile.ZipFile(path) as zf:
                names = zf.namelist()
            if any(n.startswith('word/') for n in names):
                return '.docx'
            if any(n.startswith('ppt/') for n in names):
                return '.pptx'
    except Exception as e:
        logging.debug(f"Document type sniffing failed: {e}")
    return None


async def read_document(
    file_path: str,
    extract_images: bool = False
) -> Union[str, TextContent]:
    """
    Read and extract content from documents (PDF, DOCX, PPTX).
    
    Args:
        file_path: Path to the document file (or URL)
        extract_images: Whether to extract images from document
        
    Returns:
        TextContent with extracted document content
    """
    try:
        # Handle URL downloads
        if is_url(file_path):
            logging.info(f"📥 Downloading document from URL")
            temp_path, _ = download_file_from_url(file_path)
            file_path = temp_path
        
        path = validate_file_path(file_path)
        
        logging.info(f"📄 Reading document: {path}")
        
        file_ext = path.suffix.lower()
        if file_ext not in ('.pdf', '.docx', '.pptx'):
            # The suffix came from the URL path and may be meaningless
            # (".07041", ".tmp"), so fall back to the file's magic bytes.
            file_ext = _sniff_document_type(path) or file_ext
        
        # PDF extraction
        if file_ext == '.pdf':
            import PyPDF2
            
            with open(path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                text = ""
                for page in reader.pages:
                    text += page.extract_text() + "\n"
                
                result = {
                    "file_name": path.name,
                    "file_type": "pdf",
                    "page_count": len(reader.pages),
                    "text": text[:10000],  # Limit size
                    "text_length": len(text)
                }
        
        # DOCX extraction
        elif file_ext == '.docx':
            from docx import Document
            
            doc = Document(path)
            text = "\n".join([para.text for para in doc.paragraphs])
            
            result = {
                "file_name": path.name,
                "file_type": "docx",
                "paragraph_count": len(doc.paragraphs),
                "text": text[:10000],
                "text_length": len(text)
            }
        
        # PPTX extraction
        elif file_ext == '.pptx':
            from pptx import Presentation
            
            prs = Presentation(path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            
            result = {
                "file_name": path.name,
                "file_type": "pptx",
                "slide_count": len(prs.slides),
                "text": text[:10000],
                "text_length": len(text)
            }
        
        else:
            raise ValueError(f"Unsupported file type: {file_ext}")
        
        logging.info(f"✅ Successfully extracted document content")
        
        action_response = ActionResponse(
            success=True,
            message=result,
            metadata={"file_path": str(path), "file_type": file_ext}
        )
        
        return TextContent(
            type="text",
            text=json.dumps(action_response.model_dump())
        )
        
    except Exception as e:
        error_msg = f"Document reading failed: {str(e)}"
        logging.error(f"Document error: {traceback.format_exc()}")
        
        action_response = ActionResponse(
            success=False,
            message=error_msg,
            metadata={"error_type": "document_error"}
        )
        
        return TextContent(
            type="text",
            text=json.dumps(action_response.model_dump())
        )


async def parse_image(
    image_path: str,
    use_llm: bool = True
) -> Union[str, TextContent]:
    """
    Parse and understand image content.
    
    Args:
        image_path: Path to image file or URL
        use_llm: Whether to use LLM for image understanding
        
    Returns:
        TextContent with image analysis
    """
    try:
        # Handle URL downloads
        if is_url(image_path):
            logging.info(f"📥 Downloading image from URL")
            temp_path, _ = download_file_from_url(image_path)
            image_path = temp_path
        
        path = validate_file_path(image_path)
        
        logging.info(f"🖼️ Parsing image: {path}")
        
        from PIL import Image
        
        img = Image.open(path)
        
        result = {
            "file_name": path.name,
            "format": img.format,
            "mode": img.mode,
            "size": img.size,
            "width": img.width,
            "height": img.height
        }
        
        # If LLM analysis requested, encode image for vision API
        if use_llm:
            with open(path, 'rb') as img_file:
                img_base64 = base64.b64encode(img_file.read()).decode('utf-8')
                result["base64_data"] = img_base64[:100] + "..."  # Truncated for display
                result["note"] = "Full base64 data available for vision API analysis"
        
        logging.info(f"✅ Successfully parsed image")
        
        action_response = ActionResponse(
            success=True,
            message=result,
            metadata={"file_path": str(path)}
        )
        
        return TextContent(
            type="text",
            text=json.dumps(action_response.model_dump())
        )
        
    except Exception as e:
        error_msg = f"Image parsing failed: {str(e)}"
        logging.error(f"Image error: {traceback.format_exc()}")
        
        action_response = ActionResponse(
            success=False,
            message=error_msg,
            metadata={"error_type": "image_error"}
        )
        
        return TextContent(
            type="text",
            text=json.dumps(action_response.model_dump())
        )


async def parse_video(
    video_path: str,
    extract_frames: bool = False,
    frame_interval: int = 30
) -> Union[str, TextContent]:
    """
    Parse and extract information from video files.
    
    Args:
        video_path: Path to video file or URL
        extract_frames: Whether to extract sample frames
        frame_interval: Extract one frame every N seconds
        
    Returns:
        TextContent with video metadata
    """
    try:
        # Handle URL downloads
        if is_url(video_path):
            logging.info(f"📥 Downloading video from URL")
            temp_path, _ = download_file_from_url(video_path, max_size_mb=500)
            video_path = temp_path
        
        path = validate_file_path(video_path)
        
        logging.info(f"🎥 Parsing video: {path}")
        
        import cv2
        
        video = cv2.VideoCapture(str(path))
        
        fps = video.get(cv2.CAP_PROP_FPS)
        frame_count = int(video.get(cv2.CAP_PROP_FRAME_COUNT))
        width = int(video.get(cv2.CAP_PROP_FRAME_WIDTH))
        height = int(video.get(cv2.CAP_PROP_FRAME_HEIGHT))
        duration = frame_count / fps if fps > 0 else 0
        
        result = {
            "file_name": path.name,
            "duration_seconds": duration,
            "fps": fps,
            "frame_count": frame_count,
            "resolution": f"{width}x{height}",
            "width": width,
            "height": height
        }
        
        video.release()
        
        logging.info(f"✅ Successfully parsed video metadata")
        
        action_response = ActionResponse(
            success=True,
            message=result,
            metadata={"file_path": str(path)}
        )
        
        return TextContent(
            type="text",
            text=json.dumps(action_response.model_dump())
        )
        
    except Exception as e:
        error_msg = f"Video parsing failed: {str(e)}"
        logging.error(f"Video error: {traceback.format_exc()}")
        
        action_response = ActionResponse(
            success=False,
            message=error_msg,
            metadata={"error_type": "video_error"}
        )
        
        return TextContent(
            type="text",
            text=json.dumps(action_response.model_dump())
        )


async def download_youtube_video(
    url: str,
    output_dir: str = ".",
    max_resolution: str = "720p"
) -> Union[str, TextContent]:
    """
    Download YouTube video using yt-dlp.
    
    Args:
        url: YouTube video URL
        output_dir: Directory to save video
        max_resolution: Maximum resolution (360p, 480p, 720p, 1080p)
        
    Returns:
        TextContent with download result
    """
    try:
        logging.info(f"📥 Downloading YouTube video: {url}")
        
        try:
            import yt_dlp
            
            output_template = Path(output_dir) / '%(title)s.%(ext)s'
            
            ydl_opts = {
                'format': f'best[height<={max_resolution[:-1]}]',
                'outtmpl': str(output_template),
                'quiet': False
            }
            
            with yt_dlp.YoutubeDL(ydl_opts) as ydl:
                info = ydl.extract_info(url, download=True)
                
                result = {
                    "title": info['title'],
                    "duration": info.get('duration'),
                    "output_dir": output_dir,
                    "resolution": max_resolution,
                    "video_id": info['id']
                }
                
                logging.info(f"✅ Downloaded: {info['title']}")
                
                action_response = ActionResponse(
                    success=True,
                    message=result,
                    metadata={"url": url}
                )
                
        except ImportError:
            raise ImportError("yt-dlp not installed. Install with: pip install yt-dlp")
        
        return TextContent(
            type="text",
            text=json.dumps(action_response.model_dump())
        )
        
    except Exception as e:
        error_msg = f"YouTube download failed: {str(e)}"
        logging.error(f"YouTube download error: {traceback.format_exc()}")
        
        action_response = ActionResponse(
            success=False,
            message=error_msg,
            metadata={"error_type": "youtube_download_error"}
        )
        
        return TextContent(
            type="text",
            text=json.dumps(action_response.model_dump())
        )


async def extract_youtube_transcript(
    video_id: str,
    language_code: str = "en",
    translate_to_language: str | None = None
) -> Union[str, TextContent]:
    """
    Extract transcript from a YouTube video.
    
    Args:
        video_id: YouTube video ID or URL
        language_code: Language code for the transcript (default: en)
        translate_to_language: Translate transcript to this language if provided
        
    Returns:
        TextContent with transcript data
    """
    try:
        from youtube_transcript_api import YouTubeTranscriptApi
        
        # Clean video_id if full URL was provided
        if "youtube.com" in video_id or "youtu.be" in video_id:
            if "?v=" in video_id:
                video_id = video_id.split("?v=")[-1].split("&")[0]
            elif "youtu.be/" in video_id:
                video_id = video_id.split("youtu.be/")[-1].split("?")[0]
        
        logging.info(f"📺 Extracting transcript for video ID: {video_id}")
        
        # Get transcript using correct API
        if translate_to_language:
            transcript_list = YouTubeTranscriptApi().list(video_id)
            try:
                transcript = transcript_list.find_transcript([language_code])
            except Exception:
                # If specified language not found, get any available transcript
                transcript = transcript_list.find_generated_transcript(["en"])
            # Translate to target language
            fetched_transcript = transcript.translate(translate_to_language).fetch()
            transcript_data = fetched_transcript.snippets
        else:
            try:
                # Use fetch method which returns FetchedTranscript
                fetched_transcript = YouTubeTranscriptApi().fetch(
                    video_id, 
                    languages=(language_code,)
                )
                transcript_data = fetched_transcript.snippets
            except Exception:
                # Fallback to English
                fetched_transcript = YouTubeTranscriptApi().fetch(video_id, languages=("en",))
                transcript_data = fetched_transcript.snippets
        
        # Format transcript
        formatted_transcript = []
        for entry in transcript_data:
            # Access as object attributes, not dictionary
            start_time = entry.start if hasattr(entry, 'start') else entry.get('start', 0)
            text = entry.text if hasattr(entry, 'text') else entry.get('text', '')
            minutes, seconds = divmod(int(start_time), 60)
            timestamp = f"{minutes:02d}:{seconds:02d}"
            formatted_transcript.append({
                "timestamp": timestamp,
                "text": text
            })
        
        # Create full text version
        full_text = " ".join([
            entry.text if hasattr(entry, 'text') else entry.get('text', '') 
            for entry in transcript_data
        ])
        
        result = {
            "video_id": video_id,
            "language": translate_to_language if translate_to_language else language_code,
            "transcript": formatted_transcript[:100],  # Limit to first 100 entries
            "total_entries": len(transcript_data),
            "full_text": full_text[:5000],  # Limit full text to 5000 chars
            "full_text_length": len(full_text)
        }
        
        logging.info(f"✅ Successfully extracted transcript ({len(transcript_data)} entries)")
        
        action_response = ActionResponse(
            success=True,
            message=result,
            metadata={
                "video_id": video_id,
                "language": language_code,
                "translated": translate_to_language is not None
            }
        )
        
        return TextContent(
            type="text",
            text=json.dumps(action_response.model_dump())
        )
        
    except Exception as e:
        error_msg = f"YouTube transcript extraction failed: {str(e)}"
        logging.error(f"YouTube error: {traceback.format_exc()}")
        
        action_response = ActionResponse(
            success=False,
            message=error_msg,
            metadata={"error_type": "youtube_error", "video_id": video_id}
        )
        
        return TextContent(
            type="text",
            text=json.dumps(action_response.model_dump())
        )
